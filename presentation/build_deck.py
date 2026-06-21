#!/usr/bin/env python3
"""
Generates the executive briefing deck for the endpoint-management proposal.

Audience: an executive (CIO/CFO/CEO) of the fictional Northbridge Community Bank.
Goal: explain *why* to adopt this open-source endpoint management solution — framed
around risk, cost, and audit exposure, not technical detail.

Run:  python3 presentation/build_deck.py
Out:  presentation/Northbridge_EndpointManagement_ExecBrief.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ---- Theme -----------------------------------------------------------------
NAVY = RGBColor(0x0B, 0x25, 0x45)        # primary / titles
BLUE = RGBColor(0x1B, 0x6C, 0xA8)        # accent
TEAL = RGBColor(0x13, 0xA8, 0x9E)        # secondary accent
LIGHT = RGBColor(0xF2, 0xF5, 0xF8)       # light panel
GRAY = RGBColor(0x5A, 0x6B, 0x7B)        # body subtle
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
RED = RGBColor(0xC0, 0x39, 0x2B)         # risk
GREEN = RGBColor(0x2E, 0x8B, 0x57)       # positive

FONT = "Calibri"
FONT_H = "Calibri Light"

EMU_W, EMU_H = Inches(13.333), Inches(7.5)

prs = Presentation()
prs.slide_width = EMU_W
prs.slide_height = EMU_H
BLANK = prs.slide_layouts[6]


# ---- Helpers ---------------------------------------------------------------
def slide():
    return prs.slides.add_slide(BLANK)


def rect(s, x, y, w, h, fill=None, line=None, shape=MSO_SHAPE.RECTANGLE):
    sp = s.shapes.add_shape(shape, x, y, w, h)
    sp.shadow.inherit = False
    if fill is None:
        sp.fill.background()
    else:
        sp.fill.solid()
        sp.fill.fore_color.rgb = fill
    if line is None:
        sp.line.fill.background()
    else:
        sp.line.color.rgb = line
        sp.line.width = Pt(1)
    return sp


def textbox(s, x, y, w, h, anchor=MSO_ANCHOR.TOP):
    tb = s.shapes.add_textbox(x, y, w, h)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    return tb, tf


def setpara(p, text, size, color=NAVY, bold=False, font=FONT, align=PP_ALIGN.LEFT,
            space_after=8, space_before=0):
    p.text = text
    p.alignment = align
    p.space_after = Pt(space_after)
    p.space_before = Pt(space_before)
    r = p.runs[0]
    r.font.size = Pt(size)
    r.font.bold = bold
    r.font.name = font
    r.font.color.rgb = color
    return p


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def bg(s, color=WHITE):
    rect(s, 0, 0, EMU_W, EMU_H, fill=color)


def topbar(s, kicker, title, on_dark=False):
    """Standard content-slide header: thin accent bar, kicker, title."""
    rect(s, Inches(0.6), Inches(0.55), Inches(0.12), Inches(0.85), fill=TEAL)
    tb, tf = textbox(s, Inches(0.85), Inches(0.45), Inches(11.8), Inches(1.15))
    p = tf.paragraphs[0]
    setpara(p, kicker.upper(), 12, color=(TEAL if not on_dark else TEAL), bold=True,
            space_after=2)
    p2 = tf.add_paragraph()
    setpara(p2, title, 30, color=(WHITE if on_dark else NAVY), bold=True, font=FONT_H)


def footer(s, idx, on_dark=False):
    tb, tf = textbox(s, Inches(0.6), Inches(7.0), Inches(12.1), Inches(0.4))
    p = tf.paragraphs[0]
    c = (RGBColor(0xB8, 0xC4, 0xD0) if on_dark else GRAY)
    setpara(p, "Northbridge Community Bank  ·  Endpoint Management Proposal  ·  Confidential",
            9, color=c)
    tb2, tf2 = textbox(s, Inches(12.4), Inches(7.0), Inches(0.6), Inches(0.4))
    p2 = tf2.paragraphs[0]
    setpara(p2, str(idx), 9, color=c, align=PP_ALIGN.RIGHT)


def bullets(tf, items, size=16, color=NAVY, gap=10):
    for i, (txt, lvl) in enumerate(items):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        bullet = "•   " if lvl == 0 else "–   "
        indent = "" if lvl == 0 else "      "
        setpara(p, indent + bullet + txt, size - (lvl * 1), color=color,
                space_after=gap)
        p.level = lvl


def card(s, x, y, w, h, title, body, accent=BLUE, title_size=16, body_size=12.5):
    rect(s, x, y, w, h, fill=LIGHT)
    rect(s, x, y, w, Inches(0.12), fill=accent)
    tb, tf = textbox(s, x + Inches(0.25), y + Inches(0.28), w - Inches(0.5),
                     h - Inches(0.5))
    p = tf.paragraphs[0]
    setpara(p, title, title_size, color=NAVY, bold=True, space_after=6)
    for line in body:
        pp = tf.add_paragraph()
        setpara(pp, line, body_size, color=GRAY, space_after=5)


n = 0


def num():
    global n
    n += 1
    return n


# ===========================================================================
# 1. TITLE
# ===========================================================================
s = slide()
bg(s, NAVY)
rect(s, 0, Inches(5.3), EMU_W, Inches(0.08), fill=TEAL)
rect(s, 0, Inches(5.38), EMU_W, Inches(2.12), fill=RGBColor(0x09, 0x1E, 0x39))
tb, tf = textbox(s, Inches(0.9), Inches(1.7), Inches(11.5), Inches(2.6))
setpara(tf.paragraphs[0], "Securing Every Endpoint —", 44, color=WHITE, bold=True,
        font=FONT_H, space_after=4)
setpara(tf.add_paragraph(), "Without the Enterprise Price Tag", 44, color=TEAL,
        bold=True, font=FONT_H, space_after=14)
setpara(tf.add_paragraph(),
        "An open-source approach to knowing, securing, and patching every device we own.",
        18, color=RGBColor(0xC8, 0xD4, 0xE0))
tb2, tf2 = textbox(s, Inches(0.9), Inches(5.7), Inches(11.5), Inches(1.4))
setpara(tf2.paragraphs[0], "Executive Briefing  ·  Information Technology", 15,
        color=WHITE, bold=True, space_after=4)
setpara(tf2.add_paragraph(), "Northbridge Community Bank   |   Prepared for the Executive Team",
        13, color=RGBColor(0xB8, 0xC4, 0xD0))
notes(s, "Goal of this briefing: get approval to run a low-cost, low-risk proof of "
         "concept. The headline: we can close a real security and audit gap using proven "
         "open-source software, at a fraction of commercial cost. Keep it to ~15 minutes.")

# ===========================================================================
# 2. THE QUESTION
# ===========================================================================
s = slide()
bg(s, LIGHT)
rect(s, 0, 0, Inches(0.25), EMU_H, fill=TEAL)
tb, tf = textbox(s, Inches(1.2), Inches(2.2), Inches(11), Inches(3), anchor=MSO_ANCHOR.MIDDLE)
setpara(tf.paragraphs[0], "“Is every laptop we own encrypted and fully patched?”",
        34, color=NAVY, bold=True, font=FONT_H, space_after=16)
setpara(tf.add_paragraph(),
        "Today, we answer that question with spreadsheets and guesswork — not in minutes, "
        "and not with confidence.",
        20, color=GRAY)
footer(s, num())
notes(s, "Open with the question an auditor, regulator, or the board could ask at any "
         "time. The honest answer today is 'we're not sure' — and that is the gap this "
         "proposal closes. Let it land before moving on.")

# ===========================================================================
# 3. WHERE WE STAND TODAY
# ===========================================================================
s = slide()
bg(s)
topbar(s, "Current state", "Where we stand today")
items = [
    ("~500 Windows devices — roughly 350 in-office desktops and 150 laptops that leave the building.", 0),
    ("Patching is manual and inconsistent — no guarantee a laptop in the field is up to date.", 0),
    ("No central inventory — we can't quickly say who has what, or its security status.", 0),
    ("Compliance evidence is a fire drill — assembled by hand each time an audit asks.", 0),
    ("No budget for commercial tools — Intune, Tanium, and the like are priced per device.", 0),
]
tb, tf = textbox(s, Inches(0.95), Inches(1.9), Inches(11.4), Inches(4.5))
bullets(tf, items, size=18, gap=16)
footer(s, num())
notes(s, "Plain statement of reality. No blame — the team has done well with manual "
         "process, but it doesn't scale and it doesn't give us assurance. Emphasize the "
         "roaming laptops: they're the hardest to reach and the highest risk.")

# ===========================================================================
# 4. WHAT'S AT STAKE
# ===========================================================================
s = slide()
bg(s, NAVY)
topbar(s, "The stakes", "Why this matters for a bank", on_dark=True)
x = Inches(0.85)
y = Inches(1.95)
w = Inches(3.75)
h = Inches(4.3)
gap = Inches(0.18)
risks = [
    ("Data breach", "One lost or stolen unencrypted laptop can expose customer data — a "
     "reportable breach with regulatory penalties.", RED),
    ("Ransomware / exploit", "Unpatched devices are the most common entry point for "
     "attackers. Every missing patch is an open door.", RED),
    ("Audit & regulatory findings", "Inability to prove device compliance on demand leads "
     "to findings, remediation costs, and scrutiny.", RGBColor(0xD8, 0x8C, 0x2A)),
]
for i, (t, b, ac) in enumerate(risks):
    cx = x + i * (w + gap)
    rect(s, cx, y, w, h, fill=RGBColor(0x10, 0x2E, 0x52))
    rect(s, cx, y, w, Inches(0.14), fill=ac)
    tb, tf = textbox(s, cx + Inches(0.3), y + Inches(0.45), w - Inches(0.6), h - Inches(0.8))
    setpara(tf.paragraphs[0], t, 19, color=WHITE, bold=True, space_after=12)
    setpara(tf.add_paragraph(), b, 14, color=RGBColor(0xC8, 0xD4, 0xE0))
footer(s, num(), on_dark=True)
notes(s, "Translate the technical gap into business risk an exec feels: breach liability, "
         "ransomware downtime, and regulatory findings. For a community bank, trust and "
         "reputation are the franchise — this protects them.")

# ===========================================================================
# 5. THE OPPORTUNITY
# ===========================================================================
s = slide()
bg(s, LIGHT)
rect(s, 0, 0, Inches(0.25), EMU_H, fill=TEAL)
tb, tf = textbox(s, Inches(1.2), Inches(2.1), Inches(11), Inches(3.2), anchor=MSO_ANCHOR.MIDDLE)
setpara(tf.paragraphs[0], "The opportunity", 13, color=TEAL, bold=True, space_after=10)
setpara(tf.add_paragraph(),
        "One system that knows every device, proves compliance on demand, and fixes "
        "problems automatically —",
        28, color=NAVY, bold=True, font=FONT_H, space_after=10)
setpara(tf.add_paragraph(),
        "built entirely on proven open-source software and hosted on infrastructure we "
        "already control.",
        20, color=BLUE, bold=True)
footer(s, num())
notes(s, "The pivot from problem to solution. Three verbs to remember: KNOW, PROVE, FIX. "
         "Stress 'open-source' and 'infrastructure we control' — that's the cost and "
         "data-sovereignty story execs care about.")

# ===========================================================================
# 6. WHAT IT DOES — 3 PILLARS
# ===========================================================================
s = slide()
bg(s)
topbar(s, "The solution", "Three things it does")
pillars = [
    ("SEE", "Complete visibility",
     ["Automatic inventory of every device", "Owner, hardware, OS, software",
      "Patch level and last check-in"], BLUE),
    ("PROVE", "Compliance on demand",
     ["Pass/fail against our security policy", "Encryption, antivirus, firewall, patching",
      "Audit-ready reports in minutes"], TEAL),
    ("FIX", "Automatic remediation",
     ["Push security patches and updates", "Deploy and update software",
      "Reaches laptops anywhere, securely"], GREEN),
]
x = Inches(0.85)
w = Inches(3.75)
gap = Inches(0.18)
for i, (tag, title, body, ac) in enumerate(pillars):
    cx = x + i * (w + gap)
    rect(s, cx, Inches(2.0), w, Inches(4.2), fill=LIGHT)
    rect(s, cx, Inches(2.0), w, Inches(0.7), fill=ac)
    tb, tf = textbox(s, cx, Inches(2.08), w, Inches(0.55), anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], tag, 22, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb2, tf2 = textbox(s, cx + Inches(0.3), Inches(2.95), w - Inches(0.6), Inches(3.0))
    setpara(tf2.paragraphs[0], title, 17, color=NAVY, bold=True, space_after=12)
    for line in body:
        pp = tf2.add_paragraph()
        setpara(pp, "•  " + line, 13.5, color=GRAY, space_after=8)
footer(s, num())
notes(s, "The heart of the pitch. KNOW / PROVE / FIX. Note that 'FIX' working for roaming "
         "laptops is the differentiator — the device phones home securely, so we reach it "
         "even when it's at an employee's house.")

# ===========================================================================
# 7. HOW IT WORKS (simple architecture)
# ===========================================================================
s = slide()
bg(s)
topbar(s, "How it works", "Simple, secure, and built for laptops that roam")

# Devices box
rect(s, Inches(0.9), Inches(2.4), Inches(3.2), Inches(2.6), fill=LIGHT)
rect(s, Inches(0.9), Inches(2.4), Inches(3.2), Inches(0.12), fill=BLUE)
tb, tf = textbox(s, Inches(1.05), Inches(2.7), Inches(2.9), Inches(2.2))
setpara(tf.paragraphs[0], "Our devices", 16, color=NAVY, bold=True, space_after=8)
for line in ["Desktops in the office", "Laptops anywhere", "A tiny secure agent runs quietly"]:
    setpara(tf.add_paragraph(), "•  " + line, 13, color=GRAY, space_after=7)

# Arrow 1
a1 = rect(s, Inches(4.25), Inches(3.45), Inches(1.25), Inches(0.5), fill=TEAL,
          shape=MSO_SHAPE.RIGHT_ARROW)
tb, tf = textbox(s, Inches(4.0), Inches(2.95), Inches(1.75), Inches(0.5))
setpara(tf.paragraphs[0], "Secure HTTPS", 11, color=GRAY, bold=True, align=PP_ALIGN.CENTER)

# Cloud box
rect(s, Inches(5.65), Inches(2.4), Inches(3.2), Inches(2.6), fill=NAVY)
rect(s, Inches(5.65), Inches(2.4), Inches(3.2), Inches(0.12), fill=TEAL)
tb, tf = textbox(s, Inches(5.8), Inches(2.7), Inches(2.9), Inches(2.2))
setpara(tf.paragraphs[0], "Our cloud (AWS)", 16, color=WHITE, bold=True, space_after=8)
for line in ["Central management system", "Stores inventory & compliance", "Sends patches & fixes"]:
    setpara(tf.add_paragraph(), "•  " + line, 13, color=RGBColor(0xC8, 0xD4, 0xE0), space_after=7)

# Arrow 2
rect(s, Inches(9.0), Inches(3.45), Inches(1.25), Inches(0.5), fill=TEAL,
     shape=MSO_SHAPE.RIGHT_ARROW)

# IT box
rect(s, Inches(10.4), Inches(2.4), Inches(2.4), Inches(2.6), fill=LIGHT)
rect(s, Inches(10.4), Inches(2.4), Inches(2.4), Inches(0.12), fill=GREEN)
tb, tf = textbox(s, Inches(10.55), Inches(2.7), Inches(2.1), Inches(2.2))
setpara(tf.paragraphs[0], "IT team", 16, color=NAVY, bold=True, space_after=8)
for line in ["One dashboard", "Fleet-wide view", "Act in a click"]:
    setpara(tf.add_paragraph(), "•  " + line, 13, color=GRAY, space_after=7)

tb, tf = textbox(s, Inches(0.9), Inches(5.4), Inches(11.9), Inches(1.0))
setpara(tf.paragraphs[0],
        "Devices reach out to our cloud — we never need to chase them. No office network or "
        "VPN required, so a laptop at an employee's kitchen table stays just as managed and secure.",
        15, color=NAVY)
footer(s, num())
notes(s, "Keep this non-technical. The one idea: devices 'phone home' over the same kind of "
         "secure connection as online banking, so we manage them wherever they are. No "
         "inbound access to the laptop is ever needed — that's also a security win.")

# ===========================================================================
# 8. BUILT ON TRUSTED OPEN SOURCE
# ===========================================================================
s = slide()
bg(s)
topbar(s, "Foundation", "Built on trusted, widely-used open source")
cards = [
    ("FleetDM + osquery", BLUE,
     ["Industry-standard device visibility", "Used by large security teams worldwide",
      "Powers our inventory & compliance"]),
    ("Tactical RMM", TEAL,
     ["Patch & software deployment", "Runs fixes on demand or on schedule",
      "Handles the 'fix it' half"]),
    ("We own it — no lock-in", GREEN,
     ["No per-device licensing fees", "Data stays in our cloud",
      "Free to grow at our pace"]),
]
x = Inches(0.85)
w = Inches(3.75)
gap = Inches(0.18)
for i, (t, ac, body) in enumerate(cards):
    card(s, x + i * (w + gap), Inches(2.1), w, Inches(3.9), t, body, accent=ac,
         title_size=17, body_size=13.5)
tb, tf = textbox(s, Inches(0.9), Inches(6.2), Inches(11.9), Inches(0.6))
setpara(tf.paragraphs[0],
        "These are mature tools trusted in production by thousands of organizations — not "
        "an experiment we're inventing from scratch.",
        14, color=GRAY)
footer(s, num())
notes(s, "Pre-empt the 'is open source safe/serious?' question. These are battle-tested "
         "tools. Open source here means transparency and no licensing tax — not amateur "
         "hour. We assemble proven parts rather than build from zero.")

# ===========================================================================
# 9. COMPLIANCE & SECURITY BUILT IN
# ===========================================================================
s = slide()
bg(s, NAVY)
topbar(s, "Assurance", "Compliance and security, built in", on_dark=True)
checks = [
    "Disk encryption (BitLocker) on every device",
    "Antivirus active with current definitions",
    "Firewall enabled everywhere",
    "Critical security patches applied on time",
    "Screen auto-locks when unattended",
    "Required security software present",
    "Audit-ready evidence, exported in minutes",
    "All management traffic encrypted end to end",
]
x = Inches(0.95)
y = Inches(2.05)
colw = Inches(5.8)
for i, c in enumerate(checks):
    col = i % 2
    row = i // 2
    cx = x + col * Inches(6.0)
    cy = y + row * Inches(1.0)
    rect(s, cx, cy + Inches(0.05), Inches(0.32), Inches(0.32), fill=TEAL,
         shape=MSO_SHAPE.OVAL)
    tbc, tfc = textbox(s, cx, cy + Inches(0.02), Inches(0.32), Inches(0.32),
                       anchor=MSO_ANCHOR.MIDDLE)
    setpara(tfc.paragraphs[0], "✓", 14, color=WHITE, bold=True, align=PP_ALIGN.CENTER)
    tb, tf = textbox(s, cx + Inches(0.45), cy - Inches(0.02), colw - Inches(0.5), Inches(0.9),
                     anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], c, 14.5, color=WHITE)
footer(s, num(), on_dark=True)
notes(s, "This is the slide for the audit/compliance-minded exec. Each item maps to a real "
         "control we can prove on demand. Tie back to the opening question: now we can "
         "answer 'yes, and here's the report.'")

# ===========================================================================
# 10. THE ECONOMICS
# ===========================================================================
s = slide()
bg(s)
topbar(s, "The economics", "A fraction of the commercial cost")

# Three options: commercial vs. this solution on AWS vs. this solution on-prem
options = [
    ("Commercial tools", RED, "$3–7 / device / mo", RED,
     ["~500 devices  ≈", "$18,000–42,000 / year",
      "Ongoing per-seat fees,", "rising as we grow."]),
    ("This solution — AWS", TEAL, "~$100 / month", TEAL,
     ["≈ $1,200 / year", "+ existing IT staff time",
      "No per-seat fees;", "scales to 500+ at flat cost."]),
    ("This solution — on-prem", GREEN, "$0 hosting", GREEN,
     ["Repurpose a server", "we already own",
      "No new hardware or licensing", "— just our IT staff time."]),
]
cx0, cw, cgap = Inches(0.6), Inches(3.8), Inches(0.36)
cy0, ch = Inches(2.05), Inches(3.75)
for i, (title, ac, head, hc, subs) in enumerate(options):
    cx = cx0 + i * (cw + cgap)
    rect(s, cx, cy0, cw, ch, fill=LIGHT)
    rect(s, cx, cy0, cw, Inches(0.12), fill=ac)
    tb, tf = textbox(s, cx + Inches(0.28), cy0 + Inches(0.32), cw - Inches(0.56),
                     ch - Inches(0.5))
    setpara(tf.paragraphs[0], title, 16, color=NAVY, bold=True, space_after=12)
    setpara(tf.add_paragraph(), head, 21, color=hc, bold=True, font=FONT_H, space_after=14)
    for j, sline in enumerate(subs):
        setpara(tf.add_paragraph(), sline, 13, color=GRAY,
                space_after=(10 if j == 1 else 3))

tb, tf = textbox(s, Inches(0.6), Inches(6.0), Inches(12.1), Inches(0.9))
setpara(tf.paragraphs[0],
        "Illustrative figures. We trade software licensing cost for a modest amount of our "
        "own IT time — and we keep full control of our data. Repurposing existing on-prem "
        "hardware removes even the hosting cost.",
        13, color=GRAY)
footer(s, num())
notes(s, "Three ways to pay for this. The numbers are illustrative ranges, not quotes. "
         "Option 3 is the kicker for a cost-conscious board: if we have a server with spare "
         "capacity, hosting is effectively free — we only spend a little internal time. "
         "Even on AWS, we're an order of magnitude below commercial per-seat pricing, and "
         "in every case the data stays in-house.")

# ===========================================================================
# 11. PHASED, LOW-RISK APPROACH
# ===========================================================================
s = slide()
bg(s)
topbar(s, "The plan", "A low-risk, phased rollout")
phases = [
    ("Phase 1", "Proof of Concept", "4–6 weeks", "2 test laptops. Prove the full loop: "
     "see, prove, fix. Near-zero cost.", BLUE),
    ("Phase 2", "Pilot", "1–2 months", "~50 real devices across teams. Validate at scale "
     "and refine our policies.", TEAL),
    ("Phase 3", "Fleet rollout", "Phased", "Roll out to all ~500 devices, group by group, "
     "with full reporting.", GREEN),
]
x = Inches(0.85)
w = Inches(3.75)
gap = Inches(0.18)
for i, (ph, title, dur, body, ac) in enumerate(phases):
    cx = x + i * (w + gap)
    rect(s, cx, Inches(2.1), w, Inches(3.9), fill=LIGHT)
    rect(s, cx, Inches(2.1), w, Inches(0.9), fill=ac)
    tb, tf = textbox(s, cx + Inches(0.25), Inches(2.2), w - Inches(0.5), Inches(0.75),
                     anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], ph + "  ·  " + dur, 13, color=WHITE, bold=True)
    tb2, tf2 = textbox(s, cx + Inches(0.25), Inches(3.2), w - Inches(0.5), Inches(2.6))
    setpara(tf2.paragraphs[0], title, 18, color=NAVY, bold=True, space_after=10)
    setpara(tf2.add_paragraph(), body, 14, color=GRAY)
    if i < 2:
        rect(s, cx + w - Inches(0.05), Inches(3.85), Inches(0.28), Inches(0.4), fill=ac,
             shape=MSO_SHAPE.RIGHT_ARROW)
tb, tf = textbox(s, Inches(0.9), Inches(6.2), Inches(11.9), Inches(0.6))
setpara(tf.paragraphs[0],
        "We prove value on two devices before spending real money or touching the fleet. "
        "Every phase has a clear go / no-go decision.",
        14, color=GRAY)
footer(s, num())
notes(s, "De-risk the ask. We're not betting the bank — we're asking to prove it on two "
         "machines first. Each phase gates the next, so the exec controls the spend and pace.")

# ===========================================================================
# 12. RISKS & MITIGATIONS
# ===========================================================================
s = slide()
bg(s)
topbar(s, "Honest view", "Risks — and how we manage them")
rows = [
    ("“No vendor to call”",
     "Large communities + paid support options exist; we also fully own and control the system."),
    ("Securing the central system",
     "Locked-down cloud, encrypted traffic, restricted admin access — treated as critical infrastructure."),
    ("Relies on IT staff time",
     "Phased rollout keeps effort modest; the POC quantifies the real workload before we commit."),
    ("Open-source licensing terms",
     "Reviewed for our use; one tool's license to be formally signed off before full rollout."),
]
y = Inches(2.0)
for i, (r, m) in enumerate(rows):
    cy = y + i * Inches(1.12)
    rect(s, Inches(0.9), cy, Inches(3.7), Inches(0.95), fill=NAVY)
    tb, tf = textbox(s, Inches(1.1), cy, Inches(3.3), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf.paragraphs[0], r, 14, color=WHITE, bold=True)
    rect(s, Inches(4.7), cy, Inches(7.7), Inches(0.95), fill=LIGHT)
    tb2, tf2 = textbox(s, Inches(4.95), cy, Inches(7.3), Inches(0.95), anchor=MSO_ANCHOR.MIDDLE)
    setpara(tf2.paragraphs[0], m, 13.5, color=GRAY)
footer(s, num())
notes(s, "Showing the risks builds credibility. None are show-stoppers, and the POC is "
         "designed precisely to test the open questions (effort, fit) cheaply. Address the "
         "'no vendor' worry head-on — it's the most common executive objection.")

# ===========================================================================
# 13. THE ASK
# ===========================================================================
s = slide()
bg(s, NAVY)
rect(s, 0, Inches(5.3), EMU_W, Inches(0.08), fill=TEAL)
topbar(s, "The ask", "What we're requesting today", on_dark=True)
tb, tf = textbox(s, Inches(0.95), Inches(1.95), Inches(11.5), Inches(3.2))
setpara(tf.paragraphs[0],
        "Approval to run a 4–6 week proof of concept on two test devices.",
        24, color=WHITE, bold=True, font=FONT_H, space_after=18)
for line in [
    "Incremental cost: near zero (free software, minimal cloud spend).",
    "Outcome: a live demo — a device seen, proven compliant, and patched end to end.",
    "Decision point: a clear go / no-go on a wider pilot, backed by real evidence.",
]:
    p = tf.add_paragraph()
    setpara(p, "→   " + line, 17, color=RGBColor(0xC8, 0xD4, 0xE0), space_after=14)
tb2, tf2 = textbox(s, Inches(0.95), Inches(5.6), Inches(11.5), Inches(1.2))
setpara(tf2.paragraphs[0],
        "In six weeks we can answer the question we opened with — “yes, every device is "
        "encrypted and patched, and here is the report.”",
        16, color=TEAL, bold=True)
footer(s, num(), on_dark=True)
notes(s, "The close. Make the ask concrete, small, and reversible: six weeks, two devices, "
         "near-zero cost, a clear decision afterward. Return to the opening question to "
         "bookend the story.")

# ===========================================================================
# 14. THANK YOU
# ===========================================================================
s = slide()
bg(s, NAVY)
rect(s, Inches(0.9), Inches(3.0), Inches(0.12), Inches(1.2), fill=TEAL)
tb, tf = textbox(s, Inches(1.2), Inches(2.9), Inches(10), Inches(1.6))
setpara(tf.paragraphs[0], "Thank you", 40, color=WHITE, bold=True, font=FONT_H,
        space_after=8)
setpara(tf.add_paragraph(), "Questions & discussion", 20, color=RGBColor(0xC8, 0xD4, 0xE0))
notes(s, "Invite discussion. Have the requirements and design docs ready as backup if a "
         "technical question comes up, but keep the room on the business case.")

# ---- Save ------------------------------------------------------------------
out = os.path.join(os.path.dirname(__file__),
                   "Northbridge_EndpointManagement_ExecBrief.pptx")
prs.save(out)
print("Wrote", out, "with", len(prs.slides._sldIdLst), "slides")
